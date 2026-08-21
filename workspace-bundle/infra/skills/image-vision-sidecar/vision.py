# -*- coding: utf-8 -*-
"""
image-vision-sidecar — 讓純文字主模型（DeepSeek、GLM 等）間接讀圖。
呼叫 Groq 免費 Vision API，把 PNG/JPG/PDF/PPTX/DOCX 轉成繁體中文 Markdown 描述。

用法：
    python vision.py <檔案路徑> [--mode describe|ocr] [--output <輸出檔>]
    --mode describe  詳細描述（預設）
    --mode ocr       文字抽取

輸出：預設在檔案同目錄產生 <檔名>.vision.md
"""

import argparse
import base64
import io
import os
import re
import sys

import fitz  # PyMuPDF
from PIL import Image
from pptx import Presentation
from pptx.util import Emu
from docx import Document
from groq import Groq


# 依優先順序嘗試的視覺模型（會先動態檢查 Groq 上是否存在）
VISION_MODELS = [
    "qwen/qwen3.6-27b",  # 目前 Groq 主要視覺模型
    "meta-llama/llama-4-scout-17b-16e-instruct",  # 舊版（可能已下架，保留相容）
]


def list_available_vision_models(client: Groq) -> list:
    """向 Groq 查詢實際存在的模型，回傳 VISION_MODELS 中仍在線的。"""
    try:
        available = {m.id for m in client.models.list().data}
    except Exception:
        return list(VISION_MODELS)
    return [m for m in VISION_MODELS if m in available] or list(VISION_MODELS)


def strip_think(text: str) -> str:
    """移除 Qwen 等模型的 <think> 思考過程，只留正式回答。"""
    cleaned = re.sub(r"<think>.*?</think>", "", text, flags=re.DOTALL)
    return cleaned.strip()


def get_api_key():
    """依序從環境變數與 ~/.groq_api_key 讀取 Groq API key。"""
    key = os.environ.get("GROQ_API_KEY")
    if key:
        return key.strip()
    key_path = os.path.join(os.path.expanduser("~"), ".groq_api_key")
    if os.path.exists(key_path):
        return open(key_path, encoding="utf-8").read().strip()
    raise RuntimeError(
        "找不到 Groq API key：請設定 GROQ_API_KEY 環境變數，"
        "或把 key 存在 ~/.groq_api_key"
    )


def encode_image(image_bytes: bytes, fmt: str = "PNG") -> str:
    """圖片 bytes → base64 data URI。"""
    return f"data:image/{fmt.lower()};base64,{base64.b64encode(image_bytes).decode()}"


def extract_images_from_pdf(pdf_path: str):
    """PDF → 每頁轉成 PNG 圖片，回傳 [(名稱, base64)]。"""
    images = []
    doc = fitz.open(pdf_path)
    for i, page in enumerate(doc):
        pix = page.get_pixmap(dpi=150)
        images.append((f"第{i+1}頁", encode_image(pix.tobytes("png"))))
    return images


def extract_images_from_pptx(pptx_path: str):
    """PPTX → 抽出每頁投影片中的圖片。"""
    images = []
    prs = Presentation(pptx_path)
    for si, slide in enumerate(prs.slides, 1):
        for shape in slide.shapes:
            if shape.shape_type == 13:  # PICTURE
                image = shape.image
                images.append(
                    (
                        f"slide{si} 圖片 ({shape.width // Emu(1) // 9525}px)",
                        encode_image(image.blob, image.ext),
                    )
                )
    if not images:
        # 沒有內嵌圖片時，把每頁投影片渲染成圖
        for si, slide in enumerate(prs.slides, 1):
            try:
                import pptx.tools.render  # type: ignore

                buf = io.BytesIO()
                pptx.tools.render.render_slide(slide, buf)  # type: ignore
                images.append((f"slide{si} 渲染", encode_image(buf.getvalue())))
            except Exception:
                continue
    return images


def extract_images_from_docx(docx_path: str):
    """DOCX → 抽出內嵌圖片。"""
    images = []
    doc = Document(docx_path)
    for i, rel in enumerate(doc.part.rels.values()):
        if "image" in rel.reltype:
            img = rel.target_part.blob
            images.append((f"圖片{i+1}", encode_image(img, "png")))
    return images


def extract_images(path: str):
    """依副檔名抽圖，回傳 [(名稱, base64), ...]。"""
    ext = os.path.splitext(path)[1].lower()
    if ext in (".png", ".jpg", ".jpeg", ".webp", ".bmp", ".gif"):
        img = Image.open(path)
        buf = io.BytesIO()
        img.convert("RGB").save(buf, format="PNG")
        return [(os.path.basename(path), encode_image(buf.getvalue()))]
    if ext == ".pdf":
        return extract_images_from_pdf(path)
    if ext == ".pptx":
        return extract_images_from_pptx(path)
    if ext == ".docx":
        return extract_images_from_docx(path)
    raise ValueError(f"不支援的副檔名：{ext}（支援 PNG/JPG/PDF/PPTX/DOCX）")


def build_prompt(mode: str) -> str:
    if mode == "ocr":
        return (
            "請對這張圖片做 OCR：把圖中所有文字完整擷取出來，"
            "依原本的版面順序排列。數字、單位、表格內容都要保留。"
            "如果圖中有表格，請用 Markdown 表格呈現。"
        )
    return (
        "請用繁體中文詳細描述這張圖片的內容。包含：\n"
        "- 圖片類型（照片／示意圖／表格／圖表／流程圖等）\n"
        "- 主要物件與人物\n"
        "- 可見文字（OCR）\n"
        "- 顏色與版面配置\n"
        "- 整體氛圍或傳達的資訊\n"
        "最後用「**摘要：**」給一句話總結。"
    )


def describe_image(client, model: str, b64: str, mode: str) -> str:
    resp = client.chat.completions.create(
        model=model,
        messages=[
            {
                "role": "user",
                "content": [
                    {"type": "text", "text": build_prompt(mode)},
                    {"type": "image_url", "image_url": {"url": b64}},
                ],
            }
        ],
        max_tokens=800,
    )
    return resp.choices[0].message.content or "（無回應）"


def main():
    if sys.stdout.encoding and sys.stdout.encoding.lower() not in ("utf-8", "utf8"):
        try:
            sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding="utf-8", errors="replace")
        except Exception:
            pass
    parser = argparse.ArgumentParser(description="Vision Sidecar：讓純文字模型讀圖")
    parser.add_argument("file", help="要讀取的檔案（PNG/JPG/PDF/PPTX/DOCX）")
    parser.add_argument("--mode", choices=["describe", "ocr"], default="describe")
    parser.add_argument("--output", help="輸出 Markdown 路徑（預設：<檔名>.vision.md）")
    parser.add_argument("--model", help="指定 Vision 模型（預設自動選用第一個可用）")
    args = parser.parse_args()

    path = os.path.abspath(args.file)
    if not os.path.exists(path):
        print(f"[X] 檔案不存在：{path}")
        sys.exit(1)

    print(f"[1/4] 抽取圖片中：{os.path.basename(path)}")
    images = extract_images(path)
    if not images:
        print("[!] 沒有抽出任何圖片，可能檔案內沒有圖形內容。")
        sys.exit(1)
    print(f"[OK] 抽出 {len(images)} 張圖片")

    client = Groq(api_key=get_api_key())
    model = args.model or list_available_vision_models(client)[0]
    print(f"[AI] 使用模型：{model}（模式：{args.mode}）")

    lines = [
        f"# Vision Sidecar 報告 — {os.path.basename(path)}",
        "",
        "## 來源",
        "",
        f"- 檔案：`{path}`",
        f"- 副檔名：{os.path.splitext(path)[1].lower()}",
        f"- 抽出圖片數：{len(images)}",
        f"- 描述模型：`{model}`",
        f"- 模式：`{args.mode}`",
        "",
        "## 圖片描述",
        "",
    ]

    for i, (name, b64) in enumerate(images, 1):
        print(f"  描述第 {i}/{len(images)} 張…")
        desc = strip_think(describe_image(client, model, b64, args.mode))
        lines.append(f"### 圖 {i} — {name}")
        lines.append("")
        lines.append("**描述：**")
        lines.append("")
        lines.append(desc)
        lines.append("")

    output = args.output or (os.path.splitext(path)[0] + ".vision.md")
    with open(output, "w", encoding="utf-8") as f:
        f.write("\n".join(lines))
    print(f"[OUT] 已輸出：{output}")


if __name__ == "__main__":
    main()
